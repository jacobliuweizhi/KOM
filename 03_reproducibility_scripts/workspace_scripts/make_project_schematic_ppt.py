from __future__ import annotations

import json
import math
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

ROOT = Path(__file__).resolve().parents[1]
DEPS = ROOT / "_pptx_deps"
if DEPS.exists():
    sys.path.insert(0, str(DEPS))

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "outputs"
PPTX_PATH = OUT_DIR / "project_schematic_editable.pptx"
PNG_PATH = OUT_DIR / "project_schematic_preview.png"
README_PATH = OUT_DIR / "README.md"
QA_PATH = OUT_DIR / "project_schematic_selfcheck.json"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

COLORS = {
    "ink": RGBColor(30, 45, 64),
    "muted": RGBColor(92, 105, 122),
    "blue": RGBColor(31, 86, 145),
    "blue_dark": RGBColor(20, 58, 106),
    "blue_mid": RGBColor(70, 130, 190),
    "blue_light": RGBColor(231, 241, 250),
    "blue_pale": RGBColor(244, 248, 252),
    "gray_fill": RGBColor(247, 249, 251),
    "gray_line": RGBColor(183, 194, 207),
    "line": RGBColor(80, 103, 130),
    "accent": RGBColor(38, 132, 175),
    "safe": RGBColor(48, 121, 95),
    "white": RGBColor(255, 255, 255),
}


@dataclass(frozen=True)
class Box:
    key: str
    title: str
    items: tuple[str, ...]
    x: float
    y: float
    w: float
    h: float
    fill: RGBColor = COLORS["gray_fill"]
    title_fill: RGBColor = COLORS["blue_dark"]
    emphasize: bool = False

    @property
    def center(self) -> tuple[float, float]:
        return self.x + self.w / 2, self.y + self.h / 2

    @property
    def left_mid(self) -> tuple[float, float]:
        return self.x, self.y + self.h / 2

    @property
    def right_mid(self) -> tuple[float, float]:
        return self.x + self.w, self.y + self.h / 2

    @property
    def top_mid(self) -> tuple[float, float]:
        return self.x + self.w / 2, self.y

    @property
    def bottom_mid(self) -> tuple[float, float]:
        return self.x + self.w / 2, self.y + self.h


BOXES = [
    Box(
        "input",
        "患者数据输入",
        ("基本信息", "症状与病程", "功能状态与量表评分", "影像资料（X线等）", "既往治疗", "共病与风险因素"),
        0.35,
        1.22,
        1.55,
        4.15,
        COLORS["blue_pale"],
        COLORS["blue_dark"],
    ),
    Box(
        "ai_eval",
        "AI综合评估",
        ("影像判读", "症状与功能分析", "病因重建", "风险分层", "结构化患者画像"),
        2.18,
        1.22,
        1.55,
        4.15,
        COLORS["blue_pale"],
        COLORS["blue_dark"],
    ),
    Box(
        "agents",
        "多智能体决策引擎",
        ("运动医学智能体", "营养/体重管理智能体", "心理行为智能体", "骨科综合智能体"),
        4.02,
        1.12,
        2.35,
        4.35,
        RGBColor(235, 244, 252),
        COLORS["blue_dark"],
        True,
    ),
    Box(
        "mdt",
        "MDT仲裁与方案融合",
        ("多智能体结果整合", "冲突识别", "证据分级", "安全性校验", "最终处方融合"),
        6.75,
        1.22,
        1.65,
        4.15,
        COLORS["gray_fill"],
        RGBColor(46, 80, 122),
        True,
    ),
    Box(
        "output",
        "个体化干预方案输出",
        ("运动处方", "体重/营养管理", "药物建议", "注射/手术边界建议", "随访与再评估计划"),
        8.72,
        1.22,
        1.75,
        4.15,
        COLORS["blue_pale"],
        COLORS["blue_dark"],
    ),
    Box(
        "doctor",
        "医生人机交互验证",
        ("A臂：仅患者资料", "B臂：患者资料 + AI方案", "C臂：患者资料 + AI方案 + 推理过程", "处方质量评分", "决策时间", "AI信任度", "任务负担/可接受性"),
        10.83,
        1.12,
        2.18,
        3.75,
        COLORS["gray_fill"],
        RGBColor(47, 91, 132),
    ),
    Box(
        "rag",
        "RAG证据检索与安全校验",
        ("临床指南", "随机对照试验（RCT）", "系统综述/Meta分析", "禁忌证识别", "安全性证据", "个体化适配依据"),
        4.02,
        5.78,
        4.38,
        1.28,
        RGBColor(243, 247, 250),
        RGBColor(72, 94, 119),
    ),
    Box(
        "goal",
        "最终目标",
        ("临床可审阅", "证据可追溯", "方案可执行", "安全性可校验", "支持精准、个体化KOA管理"),
        10.83,
        5.18,
        2.18,
        1.88,
        RGBColor(239, 247, 244),
        RGBColor(42, 104, 83),
        True,
    ),
]


def inch(v: float):
    return Inches(v)


def set_text_frame(text_frame, font_size: float, color=COLORS["ink"], bold=False, align=PP_ALIGN.LEFT):
    text_frame.margin_left = inch(0.08)
    text_frame.margin_right = inch(0.08)
    text_frame.margin_top = inch(0.04)
    text_frame.margin_bottom = inch(0.04)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = FONT_CN
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(slide, name: str, text: str, x: float, y: float, w: float, h: float, size: float, color=COLORS["ink"], bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    set_text_frame(tf, size, color, bold, align)
    return shape


def style_shape(shape, fill: RGBColor, line: RGBColor, width_pt: float = 1.2, radius=True):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width_pt)
    shape.line.dash_style = MSO_LINE.SOLID


def add_module(slide, box: Box):
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(box.x), inch(box.y), inch(box.w), inch(box.h))
    base.name = f"{box.key}__container"
    style_shape(base, box.fill, COLORS["gray_line"], 1.5 if box.emphasize else 1.0)

    stripe_h = 0.42
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(box.x), inch(box.y), inch(box.w), inch(stripe_h))
    stripe.name = f"{box.key}__title_band"
    style_shape(stripe, box.title_fill, box.title_fill, 0.6)

    add_textbox(
        slide,
        f"{box.key}__title_text",
        box.title,
        box.x + 0.06,
        box.y + 0.045,
        box.w - 0.12,
        stripe_h - 0.08,
        10.0 if len(box.title) <= 10 else 9.2,
        COLORS["white"],
        True,
        PP_ALIGN.CENTER,
    )

    if box.key == "agents":
        add_agent_submodules(slide, box)
    elif box.key == "doctor":
        add_doctor_submodules(slide, box)
    elif box.key == "goal":
        add_goal_items(slide, box)
    elif box.key == "rag":
        add_rag_items(slide, box)
    else:
        add_bullets(slide, box)


def add_bullets(slide, box: Box):
    top = box.y + 0.62
    row_h = (box.h - 0.82) / len(box.items)
    for idx, item in enumerate(box.items, start=1):
        y = top + (idx - 1) * row_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(box.x + 0.12), inch(y + row_h * 0.34), inch(0.06), inch(0.06))
        dot.name = f"{box.key}__bullet_dot_{idx:02d}"
        style_shape(dot, COLORS["blue_mid"], COLORS["blue_mid"], 0.3)
        add_textbox(slide, f"{box.key}__item_{idx:02d}", item, box.x + 0.24, y, box.w - 0.34, row_h * 0.9, 7.7)


def add_agent_submodules(slide, box: Box):
    labels = box.items
    gap = 0.12
    sub_w = (box.w - 0.36 - gap) / 2
    sub_h = 0.83
    start_x = box.x + 0.18
    start_y = box.y + 0.74
    fills = [RGBColor(255, 255, 255), RGBColor(247, 251, 255), RGBColor(247, 251, 255), RGBColor(255, 255, 255)]
    for i, label in enumerate(labels):
        col = i % 2
        row = i // 2
        x = start_x + col * (sub_w + gap)
        y = start_y + row * (sub_h + 0.18)
        sub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(sub_w), inch(sub_h))
        sub.name = f"{box.key}__agent_card_{i+1:02d}"
        style_shape(sub, fills[i], COLORS["blue_mid"], 1.0)
        add_textbox(slide, f"{box.key}__agent_text_{i+1:02d}", label, x + 0.05, y + 0.08, sub_w - 0.1, sub_h - 0.16, 7.6, COLORS["ink"], True, PP_ALIGN.CENTER)

    add_textbox(
        slide,
        f"{box.key}__engine_note",
        "跨学科处方生成 | 个体化约束 | 可解释推理链",
        box.x + 0.22,
        box.y + 2.88,
        box.w - 0.44,
        0.55,
        7.1,
        COLORS["muted"],
        False,
        PP_ALIGN.CENTER,
    )
    core = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(box.x + 0.32), inch(box.y + 3.42), inch(box.w - 0.64), inch(0.58))
    core.name = f"{box.key}__core_logic"
    style_shape(core, RGBColor(222, 237, 249), COLORS["blue_mid"], 1.0)
    add_textbox(slide, f"{box.key}__core_logic_text", "候选处方 + 推理依据 + 风险提示", box.x + 0.38, box.y + 3.49, box.w - 0.76, 0.42, 7.5, COLORS["blue_dark"], True, PP_ALIGN.CENTER)


def add_doctor_submodules(slide, box: Box):
    arms = box.items[:3]
    metrics = box.items[3:]
    y = box.y + 0.62
    for i, arm in enumerate(arms, start=1):
        arm_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(box.x + 0.14), inch(y), inch(box.w - 0.28), inch(0.43))
        arm_box.name = f"{box.key}__arm_card_{i:02d}"
        style_shape(arm_box, RGBColor(255, 255, 255), COLORS["gray_line"], 0.8)
        add_textbox(slide, f"{box.key}__arm_text_{i:02d}", arm, box.x + 0.2, y + 0.04, box.w - 0.4, 0.32, 6.5, COLORS["ink"], False, PP_ALIGN.CENTER)
        y += 0.5
    add_textbox(slide, f"{box.key}__metric_label", "评估终点", box.x + 0.16, y + 0.05, box.w - 0.32, 0.26, 7.3, COLORS["blue_dark"], True, PP_ALIGN.LEFT)
    y += 0.34
    for i, metric in enumerate(metrics, start=1):
        add_textbox(slide, f"{box.key}__metric_{i:02d}", f"- {metric}", box.x + 0.18, y, box.w - 0.34, 0.27, 6.7, COLORS["ink"])
        y += 0.29


def add_goal_items(slide, box: Box):
    top = box.y + 0.62
    for idx, item in enumerate(box.items, start=1):
        y = top + (idx - 1) * 0.24
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(box.x + 0.16), inch(y + 0.06), inch(0.11), inch(0.11))
        check.name = f"{box.key}__check_{idx:02d}"
        style_shape(check, COLORS["safe"], COLORS["safe"], 0.3)
        add_textbox(slide, f"{box.key}__item_{idx:02d}", item, box.x + 0.34, y, box.w - 0.48, 0.23, 7.1, COLORS["ink"], idx == len(box.items))


def add_rag_items(slide, box: Box):
    cols = 3
    cell_w = (box.w - 0.34) / cols
    cell_h = 0.28
    start_x = box.x + 0.17
    start_y = box.y + 0.55
    for i, item in enumerate(box.items):
        col = i % cols
        row = i // cols
        x = start_x + col * cell_w
        y = start_y + row * (cell_h + 0.08)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(cell_w - 0.08), inch(cell_h))
        pill.name = f"{box.key}__evidence_pill_{i+1:02d}"
        style_shape(pill, RGBColor(255, 255, 255), COLORS["gray_line"], 0.7)
        add_textbox(slide, f"{box.key}__evidence_text_{i+1:02d}", item, x + 0.03, y + 0.03, cell_w - 0.14, cell_h - 0.06, 6.5, COLORS["ink"], False, PP_ALIGN.CENTER)


def add_arrow(slide, name: str, start: tuple[float, float], end: tuple[float, float], color=COLORS["line"], width=1.55, dashed=False):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(start[0]), inch(start[1]), inch(end[0]), inch(end[1]))
    connector.name = name
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    ln = connector.line._get_or_add_ln()
    for child in list(ln):
        if child.tag.endswith("tailEnd"):
            ln.remove(child)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    ln.append(tail)
    if dashed:
        connector.line.dash_style = MSO_LINE.DASH
    return connector


def add_slide_content(prs: Presentation):
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]

    add_textbox(
        slide,
        "page__title",
        "KOM 膝骨关节炎（KOA）AI辅助全流程管理与医生人机交互验证系统",
        0.35,
        0.22,
        12.65,
        0.38,
        17,
        COLORS["ink"],
        True,
        PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "page__subtitle",
        "患者画像 → AI综合评估 → 多智能体处方生成 → MDT仲裁融合 → 个体化输出 → 医生验证闭环",
        1.15,
        0.65,
        11.1,
        0.28,
        8.5,
        COLORS["muted"],
        False,
        PP_ALIGN.CENTER,
    )

    # Soft lane labels.
    add_textbox(slide, "lane__main_flow", "主流程", 0.35, 0.94, 1.0, 0.22, 6.8, COLORS["blue_dark"], True)
    add_textbox(slide, "lane__evidence", "证据支撑与安全校验", 4.02, 5.51, 2.0, 0.22, 6.8, COLORS["muted"], True)

    for box in BOXES:
        add_module(slide, box)

    boxes = {box.key: box for box in BOXES}
    y_main = 3.28
    add_arrow(slide, "arrow__input_to_ai", boxes["input"].right_mid, (boxes["ai_eval"].x, y_main))
    add_arrow(slide, "arrow__ai_to_agents", boxes["ai_eval"].right_mid, (boxes["agents"].x, y_main), COLORS["blue"])
    add_arrow(slide, "arrow__agents_to_mdt", boxes["agents"].right_mid, (boxes["mdt"].x, y_main), COLORS["blue"])
    add_arrow(slide, "arrow__mdt_to_output", boxes["mdt"].right_mid, (boxes["output"].x, y_main))
    add_arrow(slide, "arrow__output_to_doctor", boxes["output"].right_mid, (boxes["doctor"].x, y_main))
    add_arrow(slide, "arrow__doctor_to_goal", (boxes["doctor"].x + boxes["doctor"].w / 2, boxes["doctor"].y + boxes["doctor"].h), boxes["goal"].top_mid, COLORS["safe"])

    add_arrow(slide, "arrow__rag_to_agents", (boxes["rag"].x + 1.05, boxes["rag"].y), (boxes["agents"].x + 1.05, boxes["agents"].y + boxes["agents"].h), COLORS["accent"], 1.25, dashed=True)
    add_arrow(slide, "arrow__rag_to_mdt", (boxes["rag"].x + 3.45, boxes["rag"].y), (boxes["mdt"].x + 0.8, boxes["mdt"].y + boxes["mdt"].h), COLORS["accent"], 1.25, dashed=True)
    add_arrow(slide, "arrow__followup_feedback", (11.88, 7.03), (1.1, 7.03), COLORS["muted"], 1.05, dashed=True)
    add_textbox(slide, "feedback__label", "随访反馈与持续优化", 5.28, 7.05, 2.4, 0.23, 7.0, COLORS["muted"], False, PP_ALIGN.CENTER)

    return slide


def draw_preview():
    scale = 180
    img = Image.new("RGB", (int(SLIDE_W * scale), int(SLIDE_H * scale)), "white")
    draw = ImageDraw.Draw(img)
    font_path = Path("C:/Windows/Fonts/msyh.ttc")
    if not font_path.exists():
        font_path = Path("C:/Windows/Fonts/arial.ttf")
    title_font = ImageFont.truetype(str(font_path), 31)
    subtitle_font = ImageFont.truetype(str(font_path), 16)
    module_font = ImageFont.truetype(str(font_path), 16)
    small_font = ImageFont.truetype(str(font_path), 13)
    tiny_font = ImageFont.truetype(str(font_path), 12)

    def xy(box):
        return tuple(int(v * scale) for v in box)

    draw.text(xy((SLIDE_W / 2, 0.25)), "KOM 膝骨关节炎（KOA）AI辅助全流程管理与医生人机交互验证系统", fill=(30, 45, 64), font=title_font, anchor="ma")
    draw.text(xy((SLIDE_W / 2, 0.69)), "患者画像 → AI综合评估 → 多智能体处方生成 → MDT仲裁融合 → 个体化输出 → 医生验证闭环", fill=(92, 105, 122), font=subtitle_font, anchor="ma")

    color_map = {
        COLORS["blue_pale"]: (244, 248, 252),
        COLORS["gray_fill"]: (247, 249, 251),
        RGBColor(235, 244, 252): (235, 244, 252),
        RGBColor(243, 247, 250): (243, 247, 250),
        RGBColor(239, 247, 244): (239, 247, 244),
    }

    def wrap(text: str, max_chars: int) -> list[str]:
        lines = []
        current = ""
        for ch in text:
            if len(current) >= max_chars:
                lines.append(current)
                current = ch
            else:
                current += ch
        if current:
            lines.append(current)
        return lines

    for box in BOXES:
        fill = color_map.get(box.fill, (247, 249, 251))
        draw.rounded_rectangle(xy((box.x, box.y, box.x + box.w, box.y + box.h)), radius=18, fill=fill, outline=(183, 194, 207), width=2)
        draw.rounded_rectangle(xy((box.x, box.y, box.x + box.w, box.y + 0.42)), radius=18, fill=(20, 58, 106) if box.title_fill != COLORS["safe"] else (42, 104, 83))
        draw.text(xy((box.x + box.w / 2, box.y + 0.21)), box.title, fill="white", font=module_font, anchor="mm")
        if box.key == "agents":
            for i, item in enumerate(box.items):
                col, row = i % 2, i // 2
                sx = box.x + 0.18 + col * 1.03
                sy = box.y + 0.74 + row * 1.01
                draw.rounded_rectangle(xy((sx, sy, sx + 0.98, sy + 0.83)), radius=10, fill="white", outline=(70, 130, 190), width=2)
                draw.text(xy((sx + 0.49, sy + 0.41)), "\n".join(wrap(item, 7)), fill=(30, 45, 64), font=tiny_font, anchor="mm", align="center")
            draw.text(xy((box.x + box.w / 2, box.y + 3.15)), "跨学科处方生成 | 个体化约束 | 可解释推理链", fill=(92, 105, 122), font=tiny_font, anchor="mm")
        elif box.key == "rag":
            cols = 3
            cell_w = (box.w - 0.34) / cols
            cell_h = 0.28
            start_x = box.x + 0.17
            start_y = box.y + 0.55
            for i, item in enumerate(box.items):
                col = i % cols
                row = i // cols
                sx = start_x + col * cell_w
                sy = start_y + row * (cell_h + 0.08)
                draw.rounded_rectangle(
                    xy((sx, sy, sx + cell_w - 0.08, sy + cell_h)),
                    radius=8,
                    fill="white",
                    outline=(183, 194, 207),
                    width=1,
                )
                draw.text(xy((sx + (cell_w - 0.08) / 2, sy + cell_h / 2)), "\n".join(wrap(item, 10)), fill=(30, 45, 64), font=tiny_font, anchor="mm", align="center")
        elif box.key == "goal":
            top = box.y + 0.62
            for i, item in enumerate(box.items):
                yy = top + i * 0.24
                draw.ellipse(xy((box.x + 0.16, yy + 0.08, box.x + 0.27, yy + 0.19)), fill=(48, 121, 95))
                draw.text(xy((box.x + 0.34, yy + 0.02)), "\n".join(wrap(item, 13)), fill=(30, 45, 64), font=tiny_font)
        elif box.key == "doctor":
            yy = box.y + 0.62
            for i, item in enumerate(box.items[:3]):
                draw.rounded_rectangle(
                    xy((box.x + 0.14, yy, box.x + box.w - 0.14, yy + 0.43)),
                    radius=7,
                    fill="white",
                    outline=(183, 194, 207),
                    width=1,
                )
                draw.text(xy((box.x + box.w / 2, yy + 0.21)), "\n".join(wrap(item, 15)), fill=(30, 45, 64), font=tiny_font, anchor="mm", align="center")
                yy += 0.5
            draw.text(xy((box.x + 0.18, yy + 0.12)), "评估终点", fill=(20, 58, 106), font=tiny_font)
            yy += 0.34
            for item in box.items[3:]:
                draw.text(xy((box.x + 0.2, yy)), f"- {item}", fill=(30, 45, 64), font=tiny_font)
                yy += 0.29
        else:
            max_chars = 9 if box.w < 1.8 else 13
            top = box.y + 0.62
            row_h = (box.h - 0.82) / max(1, len(box.items))
            for i, item in enumerate(box.items):
                yy = top + i * row_h
                draw.text(xy((box.x + 0.18, yy + 0.11)), "-", fill=(70, 130, 190), font=small_font)
                draw.text(xy((box.x + 0.34, yy + 0.11)), "\n".join(wrap(item, max_chars)), fill=(30, 45, 64), font=tiny_font)

    def arrow(start, end, fill=(80, 103, 130), width=3):
        draw.line(xy((*start, *end)), fill=fill, width=width)
        dx, dy = end[0] - start[0], end[1] - start[1]
        ang = math.atan2(dy, dx)
        length = 0.11
        a1 = ang + math.pi * 0.82
        a2 = ang - math.pi * 0.82
        p1 = (end[0] + length * math.cos(a1), end[1] + length * math.sin(a1))
        p2 = (end[0] + length * math.cos(a2), end[1] + length * math.sin(a2))
        draw.polygon([xy(end), xy(p1), xy(p2)], fill=fill)

    b = {box.key: box for box in BOXES}
    for s, e, c in [
        (b["input"].right_mid, (b["ai_eval"].x, 3.28), (80, 103, 130)),
        (b["ai_eval"].right_mid, (b["agents"].x, 3.28), (31, 86, 145)),
        (b["agents"].right_mid, (b["mdt"].x, 3.28), (31, 86, 145)),
        (b["mdt"].right_mid, (b["output"].x, 3.28), (80, 103, 130)),
        (b["output"].right_mid, (b["doctor"].x, 3.28), (80, 103, 130)),
        ((b["doctor"].x + b["doctor"].w / 2, b["doctor"].y + b["doctor"].h), b["goal"].top_mid, (48, 121, 95)),
    ]:
        arrow(s, e, c)

    img.save(PNG_PATH, "PNG")


def estimate_overflow(box: Box) -> bool:
    max_chars_per_line = max(7, int(box.w * 8.2))
    available_lines = max(4, int((box.h - 0.7) / 0.22))
    needed = 0
    for item in box.items:
        needed += max(1, math.ceil(len(item) / max_chars_per_line))
    return needed > available_lines + 2


def boxes_overlap(a: Box, b: Box, pad: float = 0.04) -> bool:
    return not (
        a.x + a.w + pad <= b.x
        or b.x + b.w + pad <= a.x
        or a.y + a.h + pad <= b.y
        or b.y + b.h + pad <= a.y
    )


def run_self_check() -> dict:
    checks: dict[str, object] = {}
    checks["pptx_exists"] = PPTX_PATH.exists()
    checks["pptx_size_bytes"] = PPTX_PATH.stat().st_size if PPTX_PATH.exists() else 0
    checks["preview_exists"] = PNG_PATH.exists()

    prs = Presentation(str(PPTX_PATH))
    checks["slide_count"] = len(prs.slides)
    shapes = list(prs.slides[0].shapes)
    checks["shape_count"] = len(shapes)
    checks["connector_count"] = sum(1 for s in shapes if "arrow__" in s.name)
    checks["module_container_count"] = sum(1 for s in shapes if s.name.endswith("__container"))
    checks["editable_native_objects"] = checks["shape_count"] >= 70 and checks["connector_count"] >= 8

    with zipfile.ZipFile(PPTX_PATH, "r") as zf:
        media = [name for name in zf.namelist() if name.startswith("ppt/media/")]
    checks["embedded_media_count"] = len(media)
    checks["not_whole_page_image"] = len(media) == 0

    overlaps = []
    for i, a in enumerate(BOXES):
        for b in BOXES[i + 1 :]:
            if boxes_overlap(a, b):
                overlaps.append([a.key, b.key])
    checks["module_overlaps"] = overlaps
    checks["no_module_overlap"] = len(overlaps) == 0
    checks["possible_text_overflow_modules"] = [box.key for box in BOXES if estimate_overflow(box)]
    checks["title_coverage"] = {box.key: box.title for box in BOXES}
    checks["can_open_with_python_pptx"] = True
    checks["passed"] = all(
        [
            checks["pptx_exists"],
            checks["pptx_size_bytes"] > 15000,
            checks["preview_exists"],
            checks["slide_count"] == 1,
            checks["module_container_count"] == len(BOXES),
            checks["connector_count"] >= 8,
            checks["not_whole_page_image"],
            checks["no_module_overlap"],
            not checks["possible_text_overflow_modules"],
        ]
    )
    QA_PATH.write_text(json.dumps(checks, ensure_ascii=False, indent=2), encoding="utf-8")
    return checks


def write_readme(checks: dict):
    status = "通过" if checks.get("passed") else "需要人工复核"
    text = f"""# Project Schematic Outputs

## 生成文件

- `project_schematic_editable.pptx`：核心交付文件，PowerPoint 原生可编辑对象。
- `project_schematic_preview.png`：预览图，仅用于快速查看，不嵌入 PPT。
- `project_schematic_selfcheck.json`：脚本自检记录。

## 运行命令

```powershell
cd "{ROOT}"
& "C:\\Users\\Liu\\.cache\\codex-runtimes\\codex-primary-runtime\\dependencies\\python\\python.exe" scripts\\make_project_schematic_ppt.py
```

如果改用系统 Python，请先安装依赖：

```powershell
python -m pip install python-pptx pillow
python scripts\\make_project_schematic_ppt.py
```

## PowerPoint 中如何继续编辑

打开 `project_schematic_editable.pptx` 后，直接选中任意模块框、标题文字、子模块文字、箭头或连接线即可修改。所有主模块、子模块、文本和箭头都是 PowerPoint 原生对象。预览 PNG 没有插入到 PPT 页面中。

建议打开 PowerPoint 的“选择窗格”，可以看到对象命名采用：

- `input__*`
- `ai_eval__*`
- `agents__*`
- `rag__*`
- `mdt__*`
- `output__*`
- `doctor__*`
- `goal__*`
- `arrow__*`

这些前缀方便后续批量选择和微调。

## 自检摘要

- 自检结论：{status}
- PPTX 文件大小：{checks.get("pptx_size_bytes")} bytes
- 幻灯片数量：{checks.get("slide_count")}
- 原生对象数量：{checks.get("shape_count")}
- 箭头/连接线数量：{checks.get("connector_count")}
- 主模块容器数量：{checks.get("module_container_count")}
- 嵌入媒体数量：{checks.get("embedded_media_count")}，因此不是整图图片嵌入
- 模块重叠：{checks.get("module_overlaps")}
- 潜在文字溢出模块：{checks.get("possible_text_overflow_modules")}
"""
    README_PATH.write_text(text, encoding="utf-8")


def main():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    add_slide_content(prs)
    prs.save(PPTX_PATH)
    draw_preview()
    checks = run_self_check()
    write_readme(checks)

    print(f"PPTX: {PPTX_PATH}")
    print(f"Preview: {PNG_PATH}")
    print(f"README: {README_PATH}")
    print(f"Self-check: {QA_PATH}")
    print(f"Passed: {checks['passed']}")
    if not checks["passed"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
